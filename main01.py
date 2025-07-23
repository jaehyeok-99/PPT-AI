import os
from pptx import Presentation
import requests
import json
import time  # 시간측정

# 사용할 로컬 모델 이름 (Ollama에서 pull한 모델)
#MODEL_NAME = "gemma:7b" 
#MODEL_NAME = "llama3:70b" 
MODEL_NAME = "llama3.1:8b" 

# Ollama가 컴퓨터 내부에 생성한 로컬 서버 주소
OLLAMA_API_URL = "http://localhost:11435/api/chat"

PPT_FILE_PATH = "data/3회차 강연(드론실습).pptx"
OUTPUT_FOLDER = "output"

# 프롬프트 형식 지정 (슬라이드 언급 + 간결한 TTS 스크립트용)
PROMPT_TEMPLATE = """
# 역할
당신은 PPT의 각 슬라이드 내용을 순서대로 간결하게 설명하는 AI 아나운서입니다.

# 목표
사용자가 제공한 텍스트를 기반으로, TTS가 바로 읽을 수 있는 명확하고 구조적인 설명 스크립트를 생성합니다.

# 출력 지침
1.  스크립트는 반드시 "지금부터 프레젠테이션 내용을 요약해드리겠습니다." 라는 문장으로 시작해야 합니다.
2.  각 슬라이드의 내용을 설명할 때는 "n번 슬라이드에서는..." 과 같이 반드시 슬라이드 번호를 언급해야 합니다.
3.  내용이 비슷한 여러 슬라이드는 "n번부터 m번 슬라이드까지는..." 과 같이 묶어서 설명할 수 있습니다.
4.  문장은 최대한 간결하게, 핵심 내용 위주로 구성합니다.
5.  `###` 와 같은 헤더나 글머리 기호(`*`, `-`)는 절대 사용하지 말고, 오직 하나의 완성된 줄글로만 결과물을 생성해야 합니다.
6.  "이상으로 요약을 마칩니다." 라는 문장으로 스크립트를 끝냅니다.

# 입력 텍스트
{text}

# 최종 스크립트
"""

def extract_text_from_ppt(ppt_path):
    """PPT 파일에서 텍스트를 추출하는 함수"""
    # ... (내용은 변경 없음)
    if not os.path.exists(ppt_path):
        return f"오류: '{ppt_path}' 파일을 찾을 수 없음."
    try:
        prs = Presentation(ppt_path)
        full_text = []
        for i, slide in enumerate(prs.slides, start=1):
            full_text.append(f"\n========== 슬라이드 {i} ==========\n")
            slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame.text.strip():
                                slide_text.append(cell.text_frame.text)
            full_text.append("\n".join(slide_text))
        return "\n".join(full_text)
    except Exception as e:
        return f"오류: 파일을 처리하는 중 문제가 발생\n{e}"


def summarize_with_ollama(text):
    """Ollama 로컬 API를 사용하여 텍스트를 요약하는 함수"""
    print(f" ({MODEL_NAME})를 호출하여 요약...")
    full_prompt = PROMPT_TEMPLATE.format(text=text)
    payload = {
        "model": MODEL_NAME,
        "messages": [{"role": "user", "content": full_prompt}],
        "stream": False
    }
    try:
        response = requests.post(OLLAMA_API_URL, json=payload)
        response.raise_for_status()
        response_data = response.json()
        return response_data['message']['content']
    except requests.exceptions.RequestException as e:
        return f"Ollama API 호출 중 오류가 발생. Ollama가 백그라운드에서 실행 중인지, 방화벽이 허용하는지 확인\n오류: {e}"


def save_text_to_file(text_content, original_file_path, output_dir, suffix=""):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    base_name = os.path.splitext(os.path.basename(original_file_path))[0]
    output_path = os.path.join(output_dir, f"{base_name}{suffix}.txt")
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(text_content)
    return output_path


if __name__ == "__main__":
    try:
        import requests
    except ImportError:
        print("오류: 'requests' 라이브러리가 설치되어 있지않음. 'pip install requests'를 실행.")
        exit()

    extracted_content = extract_text_from_ppt(PPT_FILE_PATH)
    
    if not extracted_content.startswith("오류:"):
        full_text_file = save_text_to_file(extracted_content, PPT_FILE_PATH, OUTPUT_FOLDER, suffix="_full")
        print(f"✅ 텍스트 추출 완료 >> {full_text_file}")
        
        start_time = time.perf_counter()

        summary_content = summarize_with_ollama(extracted_content)
        
        summary_file = save_text_to_file(summary_content, PPT_FILE_PATH, OUTPUT_FOLDER, suffix="_summary_ollama")
        print(f"✅ 요약 완료! >> {summary_file}")

        end_time = time.perf_counter()
        elapsed_time = end_time - start_time
        print(f"⏱️ 요약 처리 시간: {elapsed_time:.2f}초")
        
        print("\n--- Ollama 요약 결과 ---")
        print(summary_content)
        print("--------------------------")
    else:
        print(f"❌ {extracted_content}")