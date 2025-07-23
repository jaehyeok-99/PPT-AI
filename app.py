import os
import streamlit as st
from pptx import Presentation
import requests
import time
import pyttsx3

# --- 설정 변수 ---
MODEL_NAME = "llama3.1:8b" 
OLLAMA_API_URL = "http://localhost:11435/api/chat"
OUTPUT_FOLDER = "output"

# --- 프롬프트 ---
PROMPT_TEMPLATE = """
당신은 'PPT 요약 AI 아나운서'입니다. 지금부터 내가 제공하는 텍스트를 분석하여, TTS가 바로 읽을 수 있는 자연스러운 스크립트 하나를 생성해야 합니다.
답변할때 "스크립트를 작성하겠습니다" 이런 식으로 말하지말고 아래 규칙에 맞게만 답변해줘

스크립트는 다음과 같은 규칙을 반드시 따라야 합니다.
첫째, 발표 전체 내용을 한 문장으로 요약하며 시작하세요.
둘째, "이제 발표 내용을 상세히 설명해드리겠습니다." 라는 문장을 이어서 말하세요.
셋째, 각 슬라이드의 핵심 내용을 순서대로 자연스럽게 연결하여 설명하세요. "이어서", "다음으로"와 같은 연결어를 사용하여 전체 내용이 하나의 이야기처럼 들리도록 만들어야 합니다.
넷째, 모든 설명이 끝나면 "이상으로 요약을 마칩니다." 라는 문장으로 마무리하세요.
다섯째, 문장은 간결하게 만들고, 헤더나 글머리 기호 없이 오직 하나의 완성된 줄글로만 작성해야 합니다.

# 입력 텍스트:
{text}

# 최종 스크립트:
"""

# --- 백엔드 함수들 ---

def text_to_speech(text, output_file_path):
    """pyttsx3를 사용하여 텍스트를 음성 파일(.mp3)로 저장하는 함수"""
    try:
        engine = pyttsx3.init()
        engine.save_to_file(text, output_file_path)
        engine.runAndWait()
        return True
    except Exception as e:
        st.error(f"TTS 파일 저장 중 오류 발생: {e}")
        return False

def extract_text_from_ppt(ppt_file):
    try:
        prs = Presentation(ppt_file)
        full_text = []
        for i, slide in enumerate(prs.slides, start=1):
            full_text.append(f"\n========== 슬라이드 {i} ==========\n")
            slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame.text.strip():
                                slide_text.append(cell.text_frame.text)
            full_text.append("\n".join(slide_text))
        return "\n".join(full_text), True
    except Exception as e:
        return f"오류: 파일을 처리하는 중 문제가 발생\n{e}", False

def summarize_with_ollama(text):
    full_prompt = PROMPT_TEMPLATE.format(text=text)
    payload = {
        "model": MODEL_NAME, "messages": [{"role": "user", "content": full_prompt}], "stream": False
    }
    try:
        response = requests.post(OLLAMA_API_URL, json=payload, timeout=600)
        response.raise_for_status()
        response_data = response.json()
        return response_data['message']['content'], True
    except requests.exceptions.RequestException as e:
        return f"Ollama API 호출 중 오류 발생. Ollama가 실행 중인지 확인하세요.\n오류: {e}", False

# --- Streamlit UI 구성 ---

st.set_page_config(page_title="PPT 요약 및 TTS", layout="wide")
st.title("PPT Ai(로컬)")
st.markdown("파워포인트(.pptx) 파일을 업로드하면, AI가 내용을 요약하고 음성으로 읽어줍니다.")
st.markdown("사용 모델: llama3.1:8b")

# 파일 업로드 위젯
uploaded_file = st.file_uploader("여기에 PPT 파일을 드래그하거나 클릭하여 업로드하세요.", type=["pptx"])


if uploaded_file is not None:
    # 임시 파일 경로 대신 바로 파일 객체 사용
    file_name = uploaded_file.name
    
    with st.spinner('파일을 분석하고 요약하는 중입니다. 잠시만 기다려주세요...'):
        start_time = time.perf_counter()

        # 1. 텍스트 추출
        extracted_content, success = extract_text_from_ppt(uploaded_file)
        if not success:
            st.error(extracted_content)
        else:
            # 2. Ollama 요약
            summary_content, success = summarize_with_ollama(extracted_content)
            if not success:
                st.error(summary_content)
            else:
                # 3. 후처리
                summary_content = summary_content.replace("###", "").replace("##", "").replace("**", "").replace("*", "").replace("-", "").strip()
                
                end_time = time.perf_counter()
                elapsed_time = end_time - start_time

                st.success("요약이 완료되었습니다!")
                
                # 4. 요약 결과 출력
                st.subheader("📜 AI 요약 결과")
                st.write(summary_content)
                st.info(f"요약 처리 시간: {elapsed_time:.2f}초")

                # 5. TTS 변환 및 오디오 플레이어 출력
                if not os.path.exists(OUTPUT_FOLDER):
                    os.makedirs(OUTPUT_FOLDER)
                
                base_name = os.path.splitext(file_name)[0]
                tts_file_path = os.path.join(OUTPUT_FOLDER, f"{base_name}_summary.mp3")
                
                if text_to_speech(summary_content, tts_file_path):
                    st.subheader("🎧 음성으로 듣기")
                    st.audio(tts_file_path, format='audio/mp3')