import os
from pptx import Presentation
import requests
import time
import pyttsx3

# 사용할 로컬 모델 이름
MODEL_NAME = "llama3.1:8b" 

# Ollama 로컬 서버 주소
OLLAMA_API_URL = "http://localhost:11435/api/chat"

PPT_FILE_PATH = "data/3회차 강연(드론실습).pptx"
OUTPUT_FOLDER = "output"

# 프롬프트 형식 지정 
PROMPT_TEMPLATE = """
너는 PPT의 각 슬라이드 내용을 순서대로 간결하게 설명하는 ppt ai 아나운서이다.

# 목표
사용자가 제공한 출력 지침에 따라 TTS가 바로 읽을 수 있는 명확하고 구조적인 설명 스크립트를 생성한다

# 출력 지침
1. 먼저, 프레젠테이션 전체 내용을 한 문장으로 요약하며 시작한다.
2. 그 다음, "이제 각 슬라이드의 내용을 상세히 요약해드리겠습니다." 라는 고정된 문장을 말한다.
3. 이어서, "n번 슬라이드에서는..." 과 같이 슬라이드 번호를 언급하며 각 슬라이드의 핵심 내용을 순서대로 설명한다.
4. 내용이 비슷한 여러 슬라이드는 "n번부터 m번 슬라이드까지는..." 과 같이 묶어서 설명할 수 있다
5. 문장은 최대한 간결하게, 핵심 내용 위주로 구성한다.
6. 모든 설명이 끝나면, "이상으로 요약을 마칩니다." 라는 문장으로 스크립트를 끝낸다

# 입력 텍스트
{text}

# 최종 스크립트
"""

def text_to_speech(text):
    """pyttsx3를 사용하여 텍스트를 음성으로 직접 출력하는 함수"""
    try:
        print("TTS 출력을 시작합니다...")
        engine = pyttsx3.init()
        engine.say(text)
        engine.runAndWait()
        print("✅ TTS 출력 완료!")
        return True
    except Exception as e:
        print(f"❌ TTS 출력 중 오류 발생: {e}")
        return False

def extract_text_from_ppt(ppt_path):
    if not os.path.exists(ppt_path):
        return f"오류: '{ppt_path}' 파일을 찾을 수 없음."
    try:
        prs = Presentation(ppt_path)
        full_text = []
        for i, slide in enumerate(prs.slides, start=1):
            full_text.append(f"\n========== 슬라이드 {i} ==========\n")
            slide_text = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            for shape in slide.shapes:
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame.text.strip():
                                slide_text.append(cell.text_frame.text)
            full_text.append("\n".join(slide_text))
        return "\n".join(full_text)
    except Exception as e:
        return f"오류: 파일을 처리하는 중 문제가 발생\n{e}"

def summarize_with_ollama(text):
    print(f"Ollama ({MODEL_NAME})를 호출하여 요약을 시작합니다...")
    full_prompt = PROMPT_TEMPLATE.format(text=text)
    payload = {
        "model": MODEL_NAME,
        "temperature": 0.9,
        "top_p": 0.9,
        "messages": [{"role": "user", "content": full_prompt}],
        "stream": False
    }
    try:
        response = requests.post(OLLAMA_API_URL, json=payload)
        response.raise_for_status()
        response_data = response.json()
        return response_data['message']['content']
    except requests.exceptions.RequestException as e:
        return f"Ollama API 호출 중 오류가 발생. Ollama가 백그라운드에서 실행 중인지, 방화벽이 허용하는지 확인\n오류: {e}"


def save_text_to_file(text_content, original_file_path, output_dir, suffix=""):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    base_name = os.path.splitext(os.path.basename(original_file_path))[0]
    output_path = os.path.join(output_dir, f"{base_name}{suffix}.txt")
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(text_content)
    return output_path


if __name__ == "__main__":
    extracted_content = extract_text_from_ppt(PPT_FILE_PATH)
    
    if not extracted_content.startswith("오류:"):
        full_text_file = save_text_to_file(extracted_content, PPT_FILE_PATH, OUTPUT_FOLDER, suffix="_full")
        print(f"✅ 텍스트 추출 완료! >> {full_text_file}")
        
        start_time = time.perf_counter()

        summary_content = summarize_with_ollama(extracted_content)

        summary_content = summary_content.replace("###", "").replace("*", "").replace("-", "").strip()
        
        end_time = time.perf_counter()
        elapsed_time = end_time - start_time
        
        summary_file = save_text_to_file(summary_content, PPT_FILE_PATH, OUTPUT_FOLDER, suffix="_summary_ollama")
        print(f"✅ 요약 완료! >> {summary_file}")
        print("\n--- Ollama 요약 결과 ---")
        print(summary_content)
        print("--------------------------")
        if summary_content and not summary_content.startswith("Ollama API 호출 중 오류가 발생"):
            text_to_speech(summary_content)
        print(f"⏱️ 처리 시간: {elapsed_time:.2f}초")
        
    else:
        print(f"❌ {extracted_content}")